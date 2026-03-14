from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


OUTPUT_FILE = "codex_enterprise_rollout.pptx"

# Theme constants
ACCENT_BLUE = RGBColor(0x12, 0x3B, 0x73)  # subtle deep blue
TEXT_DARK = RGBColor(0x1A, 0x1A, 0x1A)
SUBTLE_GRAY = RGBColor(0x63, 0x63, 0x63)
TITLE_FONT = "Aptos Display"
BODY_FONT = "Aptos"


SLIDES = [
    {
        "title": "Deploying Codex (Local) for Enterprise Engineering",
        "subtitle": "AI-assisted development for a Fortune 500 retail platform",
        "bullets": [
            "Increase developer velocity without weakening security controls",
            "Align IT risk ownership with practical engineering enablement",
            "Start with a low-risk, high-impact workflow",
            "Target measurable outcomes in onboarding, quality, and delivery",
        ],
        "visual": "3-column visual: Productivity | Governance | Business Outcomes",
    },
    {
        "title": "Current Engineering Challenges",
        "bullets": [
            "Large legacy codebase slows feature delivery",
            "Slow onboarding for new and rotating engineers",
            "Inconsistent documentation across services",
            "High time cost to trace code paths and dependencies",
            "Frequent context switching increases rework",
        ],
        "visual": "Simple bar chart: Code understanding is the largest time segment",
    },
    {
        "title": "Where Codex Fits in the Developer Workflow",
        "bullets": [
            "Accelerates code understanding in unfamiliar modules",
            "Maps dependencies and service interactions quickly",
            "Generates first-draft technical documentation",
            "Assists debugging with likely failure-path summaries",
            "Improves manager-led onboarding and mentoring efficiency",
        ],
        "visual": "Flow: Read Code -> Ask Codex -> Validate -> Implement",
    },
    {
        "title": "First Workflow: Code Understanding",
        "bullets": [
            "Highest immediate productivity gain across teams",
            "Low operational risk: read-focused, no autonomous prod actions",
            "Shortens onboarding time on legacy systems",
            "Improves design and change decisions before coding",
            "Establishes a foundation for future AI-assisted workflows",
        ],
        "visual": "2x2 matrix: Code Understanding in High Impact / Low Risk",
    },
    {
        "title": "Enterprise Guardrails",
        "bullets": [
            "Repository access control mapped to IAM/RBAC",
            "Prompt filtering for secrets and restricted data",
            "Usage telemetry for auditability and adoption tracking",
            "Policy enforcement through approved harness and sandboxing",
            "Human-in-the-loop review remains required for merges",
        ],
        "visual": "Layered architecture: Developer -> Harness -> Policy -> Repos",
    },
    {
        "title": "Live Demo Example: Inventory Availability Service",
        "bullets": [
            "Scenario: engineer investigates stock-availability logic before a change",
            'Prompt: "Explain how inventory availability is calculated in this service."',
            "Codex returns key files, calculation path, and dependencies",
            "Expected output: assumptions, edge cases, and validation hints",
            "Result: faster root-cause analysis and safer implementation",
        ],
        "visual": "Split pane mockup: Prompt on left, explanation + file map on right",
    },
    {
        "title": "Rollout Strategy and Business Impact",
        "bullets": [
            "Pilot: 2-3 teams, baseline metrics, guardrail validation",
            "Standardized adoption: playbooks, prompt patterns, policy templates",
            "Enterprise scale: governance dashboards and broad enablement",
            "Outcomes: faster onboarding, higher productivity, improved code quality",
            "KPIs: time-to-first-PR, cycle time, escaped defects, doc coverage",
        ],
        "visual": "3-stage roadmap: Pilot -> Standardized Adoption -> Enterprise Scale",
    },
]


def add_accent_bar(slide):
    # Thin top accent line for consistent enterprise visual identity.
    bar = slide.shapes.add_shape(
        autoshape_type_id=1,  # msoShapeRectangle
        left=Inches(0),
        top=Inches(0),
        width=Inches(13.333),
        height=Inches(0.12),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_BLUE
    bar.line.fill.background()


def style_title(title_shape):
    tf = title_shape.text_frame
    p = tf.paragraphs[0]
    for run in p.runs:
        run.font.name = TITLE_FONT
        run.font.size = Pt(38)
        run.font.bold = True
        run.font.color.rgb = TEXT_DARK


def add_subtitle(slide, text):
    subtitle = slide.shapes.add_textbox(Inches(0.9), Inches(1.75), Inches(11.5), Inches(0.7))
    tf = subtitle.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = BODY_FONT
    run.font.size = Pt(21)
    run.font.color.rgb = ACCENT_BLUE


def style_content_text(content_shape, bullets, visual):
    tf = content_shape.text_frame
    tf.clear()

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.space_after = Pt(8)
        p.line_spacing = 1.2
        if p.runs:
            for run in p.runs:
                run.font.name = BODY_FONT
                run.font.size = Pt(22)
                run.font.color.rgb = TEXT_DARK

    p = tf.add_paragraph()
    p.text = f"Visual suggestion: {visual}"
    p.level = 0
    p.space_before = Pt(12)
    p.space_after = Pt(0)
    if p.runs:
        for run in p.runs:
            run.font.name = BODY_FONT
            run.font.size = Pt(16)
            run.font.italic = True
            run.font.color.rgb = SUBTLE_GRAY


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_and_content = prs.slide_layouts[1]

    for idx, data in enumerate(SLIDES):
        slide = prs.slides.add_slide(title_and_content)
        add_accent_bar(slide)

        slide.shapes.title.text = data["title"]
        style_title(slide.shapes.title)

        content = slide.placeholders[1]

        if idx == 0:
            add_subtitle(slide, data["subtitle"])
            # Move content box down to avoid crowding the subtitle.
            content.left = Inches(0.95)
            content.top = Inches(2.55)
            content.width = Inches(11.8)
            content.height = Inches(3.9)
        else:
            content.left = Inches(0.95)
            content.top = Inches(1.9)
            content.width = Inches(11.8)
            content.height = Inches(4.9)

        style_content_text(content, data["bullets"], data["visual"])

    prs.save(OUTPUT_FILE)


if __name__ == "__main__":
    build_presentation()
    print(f"Created {OUTPUT_FILE}")
